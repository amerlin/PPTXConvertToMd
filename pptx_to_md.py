from __future__ import annotations

import argparse
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass(order=True)
class PositionedBlock:
    top: int
    left: int
    markdown: str


def _normalize_text(value: str) -> str:
    return value.replace("\xa0", " ").strip()


def _escape_table_cell(value: str) -> str:
    return value.replace("|", r"\|").replace("\n", "<br>").strip()


def _format_run(run) -> str:
    text = run.text.replace("\x0b", "\n")
    if not text:
        return ""

    if run.font.bold and run.font.italic:
        return f"***{text}***"
    if run.font.bold:
        return f"**{text}**"
    if run.font.italic:
        return f"*{text}*"
    return text


def _paragraph_text(paragraph) -> str:
    raw_text = "".join(_format_run(run) for run in paragraph.runs) or paragraph.text
    return _normalize_text(raw_text)


def _text_frame_to_markdown(text_frame) -> str:
    paragraphs: list[tuple[int, str]] = []
    for paragraph in text_frame.paragraphs:
        text = _paragraph_text(paragraph)
        if text:
            paragraphs.append((paragraph.level, text))

    if not paragraphs:
        return ""

    if len(paragraphs) == 1 and paragraphs[0][0] == 0:
        return paragraphs[0][1]

    return "\n".join(f"{'  ' * level}- {text}" for level, text in paragraphs)


def _table_to_markdown(table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [_escape_table_cell(_normalize_text(cell.text)) for cell in row.cells]
        if any(cells):
            rows.append(cells)

    if not rows:
        return ""

    column_count = max(len(row) for row in rows)
    normalized_rows = [row + [""] * (column_count - len(row)) for row in rows]

    header = normalized_rows[0]
    separator = ["---"] * column_count
    body = normalized_rows[1:]

    lines = [
        f"| {' | '.join(header)} |",
        f"| {' | '.join(separator)} |",
    ]
    lines.extend(f"| {' | '.join(row)} |" for row in body)
    return "\n".join(lines)


def _iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
            continue
        yield shape


def _shape_image(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return shape.image

    if getattr(shape, "is_placeholder", False):
        try:
            return shape.image
        except (AttributeError, ValueError):
            return None

    return None


def _export_shape_image(shape, slide_number: int, image_number: int, images_dir: Path) -> Path | None:
    image = _shape_image(shape)
    if image is None:
        return None

    images_dir.mkdir(parents=True, exist_ok=True)
    extension = image.ext.lower()
    file_name = f"slide_{slide_number:02d}_image_{image_number:02d}.{extension}"
    image_path = images_dir / file_name
    image_path.write_bytes(image.blob)
    return image_path


def _slide_title(slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is None or not title_shape.has_text_frame:
        return ""
    return _text_frame_to_markdown(title_shape.text_frame)


def _notes_markdown(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return _text_frame_to_markdown(slide.notes_slide.notes_text_frame)


def _slide_content(slide) -> list[PositionedBlock]:
    title_shape = slide.shapes.title
    title_shape_id = getattr(title_shape, "shape_id", None)
    blocks: list[PositionedBlock] = []

    for shape in _iter_shapes(slide.shapes):
        if getattr(shape, "shape_id", None) == title_shape_id:
            continue

        if getattr(shape, "has_text_frame", False):
            markdown = _text_frame_to_markdown(shape.text_frame)
            if markdown:
                blocks.append(PositionedBlock(shape.top, shape.left, markdown))

        if getattr(shape, "has_table", False):
            markdown = _table_to_markdown(shape.table)
            if markdown:
                blocks.append(PositionedBlock(shape.top, shape.left, markdown))

    return sorted(blocks)


def _slide_images(slide, slide_number: int, markdown_path: Path, images_dir: Path) -> list[PositionedBlock]:
    images: list[PositionedBlock] = []
    image_number = 1

    for shape in _iter_shapes(slide.shapes):
        image_path = _export_shape_image(shape, slide_number, image_number, images_dir)
        if image_path is None:
            continue

        try:
            markdown_link = Path(os.path.relpath(image_path, markdown_path.parent)).as_posix()
        except ValueError:
            markdown_link = image_path.as_uri()
        alt_text = f"Slide {slide_number} image {image_number}"
        images.append(PositionedBlock(shape.top, shape.left, f"![{alt_text}]({markdown_link})"))
        image_number += 1

    return sorted(images)


def _default_output_paths(input_path: Path, output_path: Path | None, images_dir: Path | None) -> tuple[Path, Path]:
    resolved_output = output_path or input_path.with_suffix(".md")
    resolved_images_dir = images_dir or resolved_output.with_name(f"{resolved_output.stem}_assets")
    return resolved_output, resolved_images_dir


def convert_pptx_to_markdown(
    input_path: Path,
    output_path: Path | None = None,
    images_dir: Path | None = None,
    include_notes: bool = True,
) -> Path:
    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")
    if input_path.suffix.lower() != ".pptx":
        raise ValueError("Input file must have a .pptx extension.")

    markdown_path, resolved_images_dir = _default_output_paths(input_path, output_path, images_dir)
    markdown_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(input_path))
    document_title = _normalize_text(presentation.core_properties.title or "") or input_path.stem

    lines = [
        f"# {document_title}",
        "",
        f"Source: `{input_path.name}`",
    ]

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = _slide_title(slide) or f"Slide {slide_index}"
        content_blocks = _slide_content(slide)
        image_blocks = _slide_images(slide, slide_index, markdown_path, resolved_images_dir)
        notes_markdown = _notes_markdown(slide) if include_notes else ""

        lines.extend(
            [
                "",
                f"## Slide {slide_index} - {title}",
                "",
            ]
        )

        if content_blocks:
            lines.extend(["### Content", ""])
            for block in content_blocks:
                lines.extend([block.markdown, ""])

        if image_blocks:
            lines.extend(["### Images", ""])
            for block in image_blocks:
                lines.extend([block.markdown, ""])

        if notes_markdown:
            lines.extend(["### Speaker notes", "", notes_markdown, ""])

        if not content_blocks and not image_blocks and not notes_markdown:
            lines.extend(["_No content extracted._", ""])

    markdown_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
    return markdown_path


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Convert a .pptx presentation into a Markdown document with exported slide images."
    )
    parser.add_argument("input", help="Path to the input .pptx file")
    parser.add_argument(
        "-o",
        "--output",
        help="Path to the output .md file (default: same name as input)",
    )
    parser.add_argument(
        "--images-dir",
        help="Directory where extracted images will be written (default: <output_stem>_assets)",
    )
    parser.add_argument(
        "--no-notes",
        action="store_true",
        help="Exclude speaker notes from the output",
    )
    args = parser.parse_args()

    input_path = Path(args.input).expanduser().resolve()
    output_path = Path(args.output).expanduser().resolve() if args.output else None
    images_dir = Path(args.images_dir).expanduser().resolve() if args.images_dir else None

    markdown_path = convert_pptx_to_markdown(input_path, output_path, images_dir, include_notes=not args.no_notes)
    print(f"Markdown created: {markdown_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
